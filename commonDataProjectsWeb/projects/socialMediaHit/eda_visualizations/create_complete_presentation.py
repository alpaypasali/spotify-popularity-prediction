#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Complete Final Presentation with EDA and Speechiness Analysis
"""
import sys
import os
sys.path.insert(0, os.path.dirname(__file__))

# Import from the previous script's functions
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path

# Colors
SPOTIFY_GREEN = RGBColor(30, 215, 96)
DARK_GRAY = RGBColor(18, 18, 18)
WHITE = RGBColor(255, 255, 255)

def add_title_slide(prs, title, subtitle, author, date):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_GRAY
    
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(9)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(54)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = SPOTIFY_GREEN
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    top = Inches(3.8)
    height = Inches(0.8)
    txBox2 = slide.shapes.add_textbox(left, top, width, height)
    tf2 = txBox2.text_frame
    tf2.text = subtitle
    tf2.paragraphs[0].font.size = Pt(32)
    tf2.paragraphs[0].font.color.rgb = WHITE
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    top = Inches(5.5)
    height = Inches(1)
    txBox3 = slide.shapes.add_textbox(left, top, width, height)
    tf3 = txBox3.text_frame
    tf3.text = f"{author}\n{date}"
    tf3.paragraphs[0].font.size = Pt(20)
    tf3.paragraphs[0].font.color.rgb = WHITE
    tf3.paragraphs[0].alignment = PP_ALIGN.CENTER
    if len(tf3.paragraphs) > 1:
        tf3.paragraphs[1].font.size = Pt(18)
        tf3.paragraphs[1].font.color.rgb = WHITE
        tf3.paragraphs[1].alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_GRAY
    
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = title.upper()
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = SPOTIFY_GREEN
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)
    txBox2 = slide.shapes.add_textbox(left, top, width, height)
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i > 0:
            p = tf2.add_paragraph()
        else:
            p = tf2.paragraphs[0]
        p.text = f"• {item}" if not item.startswith("•") else item
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.space_after = Pt(12)

def add_two_column_slide(prs, title, left_items, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_GRAY
    
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = title.upper()
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = SPOTIFY_GREEN
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    left_col = Inches(0.5)
    top = Inches(1.5)
    width = Inches(4.2)
    height = Inches(5)
    txBox_left = slide.shapes.add_textbox(left_col, top, width, height)
    tf_left = txBox_left.text_frame
    tf_left.word_wrap = True
    
    for i, item in enumerate(left_items):
        if i > 0:
            p = tf_left.add_paragraph()
        else:
            p = tf_left.paragraphs[0]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(10)
    
    right_col = Inches(5.3)
    txBox_right = slide.shapes.add_textbox(right_col, top, width, height)
    tf_right = txBox_right.text_frame
    tf_right.word_wrap = True
    
    for i, item in enumerate(right_items):
        if i > 0:
            p = tf_right.add_paragraph()
        else:
            p = tf_right.paragraphs[0]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(10)

def add_image_slide(prs, title, image_path, description_items=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_GRAY
    
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = title.upper()
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = SPOTIFY_GREEN
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    img_path = Path(image_path)
    if img_path.exists():
        left_img = Inches(0.5)
        top_img = Inches(1.3)
        width_img = Inches(9)
        height_img = Inches(5)
        slide.shapes.add_picture(str(img_path), left_img, top_img, width_img, height_img)
    
    if description_items:
        left_desc = Inches(0.5)
        top_desc = Inches(6.5)
        width_desc = Inches(9)
        height_desc = Inches(0.6)
        txBox_desc = slide.shapes.add_textbox(left_desc, top_desc, width_desc, height_desc)
        tf_desc = txBox_desc.text_frame
        for i, desc in enumerate(description_items):
            if i > 0:
                p = tf_desc.add_paragraph()
            else:
                p = tf_desc.paragraphs[0]
            p.text = f"• {desc}"
            p.font.size = Pt(14)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

def add_two_column_image_slide(prs, title, image_path, content_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_GRAY
    
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = title.upper()
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = SPOTIFY_GREEN
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    img_path = Path(image_path)
    if img_path.exists():
        left_img = Inches(0.5)
        top_img = Inches(1.3)
        width_img = Inches(4.5)
        height_img = Inches(5.5)
        slide.shapes.add_picture(str(img_path), left_img, top_img, width_img, height_img)
    
    left_content = Inches(5.2)
    top_content = Inches(1.3)
    width_content = Inches(4.3)
    height_content = Inches(5.5)
    txBox_content = slide.shapes.add_textbox(left_content, top_content, width_content, height_content)
    tf_content = txBox_content.text_frame
    tf_content.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i > 0:
            p = tf_content.add_paragraph()
        else:
            p = tf_content.paragraphs[0]
        p.text = f"• {item}" if not item.startswith("•") else item
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

eda_dir = Path("projects/socialMediaHit/eda_visualizations")

# 1. Kapak
add_title_slide(prs, "SOCIAL MEDIA HIT PREDICTION SUITE", 
                "MIUUL Data Scientist Bootcamp\nFinal Project",
                "Enes TOP", "Aralık 2025")

# 2-6. Existing slides (same as before)
add_content_slide(prs, "PROJE ÖZETİ", [
    "AI destekli sosyal medya içerik hit tahmin platformu",
    "Spotify müzik hit potansiyeli analizi",
    "Instagram içerik performans tahmini (Görsel + Metin)",
    "Multi-modal Reels içerik üretimi için entegre sistem",
    "Gerçek zamanlı tahmin ve öneriler",
    "Production-ready Django web uygulaması"
])

add_content_slide(prs, "PROBLEM TANIMI", [
    "Sosyal medya içerik üreticileri hangi içeriğin viral olacağını bilmiyor",
    "Manuel deneme-yanılma süreci zaman ve kaynak israfı",
    "Müzik prodüktörleri ve DJ'ler hit potansiyelini önceden göremiyor",
    "Instagram Reels içerikleri için multi-modal analiz eksikliği",
    "Veri bilimi tekniklerini pratik bir problemde uygulama ihtiyacı"
])

add_content_slide(prs, "ÇÖZÜM YAKLAŞIMI", [
    "Spotify API ve Kaggle datasets ile zengin veri toplama",
    "Machine Learning modelleri ile hit potansiyeli tahmini",
    "Multi-modal analiz: Ses + Görsel + Metin",
    "Ensemble learning (XGBoost, LightGBM, Random Forest, GBM)",
    "Real-time web arayüzü ile kullanıcı dostu deneyim",
    "Market-specific tahminler (US, GB, DE, FR, JP, BR)"
])

# 5. Sistem Mimarisi - Detaylı Akış Diyagramı
system_arch_img = eda_dir / 'system_architecture_diagram.png'
if system_arch_img.exists():
    add_image_slide(prs, "SİSTEM MİMARİSİ - HIT TAHMİN PLATFORMU", 
                   str(system_arch_img),
                   ["Spotify API ve Kaggle Datasets veri kaynakları",
                    "Multi-modal feature extraction (NLP, CNN, MFCC, Video Analytics)",
                    "6 farklı hit tahmin modülü",
                    "Sistem entegrasyonu ve API katmanı",
                    "Reels içerik hit tahmini ve üretici"])
else:
    add_two_column_slide(prs, "SİSTEM MİMARİSİ",
        ["FRONTEND", "• Django Templates", "• Tailwind CSS", "• Chart.js Visualizations", "• Vanilla JavaScript", "", "BACKEND", "• Django REST Framework", "• ML Model Services", "• API Endpoints"],
        ["DATA PROCESSING", "• Pandas & NumPy", "• Feature Engineering", "• Data Preprocessing", "", "ML MODELS", "• XGBoost Regressor", "• LightGBM", "• Random Forest", "• Gradient Boosting", "", "DEPLOYMENT", "• Docker & Docker Compose", "• PostgreSQL", "• AWS Ready"]
    )

add_two_column_slide(prs, "TEKNOLOJİ STACK",
    ["BACKEND", "• Django 4.2", "• Django REST Framework", "• Python 3.11", "", "MACHINE LEARNING", "• XGBoost", "• LightGBM", "• scikit-learn", "• Pandas & NumPy", "", "DATA SOURCES", "• Spotify Web API", "• Kaggle Datasets"],
    ["FRONTEND", "• Tailwind CSS", "• Chart.js", "• Vanilla JavaScript", "", "INFRASTRUCTURE", "• Docker", "• PostgreSQL", "• Prometheus & Grafana", "", "DEPLOYMENT", "• AWS Free Tier", "• GitHub Actions CI/CD"]
)

# 7-14. EDA Slides
add_content_slide(prs, "EXPLORATORY DATA ANALYSIS", [
    "Toplam 141,190 şarkı verisi analiz edildi",
    "17 özellik (feature) kolonu: audio features + metadata",
    "Audio features: danceability, energy, valence, loudness, tempo, speechiness",
    "Metadata: popularity, genre, release_year, explicit, emotion",
    "Çoklu veri kaynağı: Spotify API + Kaggle datasets",
    "Temporal coverage: 1960-2024 arası şarkılar"
])

# Add EDA visualizations (8-14)
img_files = [
    ("1_audio_features_distribution.png", "AUDIO FEATURES DAĞILIMI", 
     ["Tüm audio feature'lar normal dağılıma yakın", "Danceability ve energy yüksek varyans gösteriyor", "Loudness negatif değerler (dB cinsinden)"]),
    ("2_correlation_heatmap.png", "FEATURE CORRELATION MATRIX", None, True),
    ("3_feature_vs_popularity.png", "FEATURE VS POPULARITY ANALİZİ",
     ["Energy ve danceability popularity ile pozitif trend", "Acousticness düşük popularity ile ilişkili", "Valence (mood) orta seviye popularity ile ilişkili"]),
    ("4_genre_distribution.png", "GENRE DAĞILIMI", None, True),
    ("5_year_distribution.png", "YIL BAZLI DAĞILIM",
     ["Veri seti 1960-2024 arası şarkıları kapsıyor", "2010 sonrası veri yoğunluğu artıyor"]),
    ("6_feature_evolution.png", "AUDIO FEATURES EVOLUTION", None, True),
    ("7_boxplots.png", "OUTLIER ANALİZİ - BOX PLOTS",
     ["Loudness ve tempo'da önemli outlier'lar var", "Speechiness ve acousticness geniş dağılım gösteriyor", "Outlier'lar model eğitimi öncesi temizlendi"])
]

correlation_content = [
    "Energy ve loudness güçlü pozitif korelasyon (0.75+)",
    "Acousticness ve energy negatif korelasyon (-0.4)",
    "Danceability ve valence orta pozitif korelasyon (0.5)",
    "Tempo diğer feature'larla zayıf korelasyon",
    "Popularity ile energy ve danceability pozitif ilişki"
]

genre_content = [
    "Pop, rock, hip-hop en yaygın türler",
    "Electronic ve alternative önemli paya sahip",
    "Genre bazlı hit potansiyeli farklılık gösteriyor",
    "Market-specific genre tercihleri mevcut"
]

evolution_content = [
    "Danceability zamanla artış trendi gösteriyor",
    "Energy seviyeleri 2000'lerden sonra yükselmiş",
    "Tempo modern müzikte daha yüksek",
    "Müzik tarzları zamanla değişiyor"
]

for img_file, title, desc, is_two_col in [
    ("1_audio_features_distribution.png", "AUDIO FEATURES DAĞILIMI", 
     ["Tüm audio feature'lar normal dağılıma yakın", "Danceability ve energy yüksek varyans gösteriyor", "Loudness negatif değerler (dB cinsinden)"], False),
    ("2_correlation_heatmap.png", "FEATURE CORRELATION MATRIX", correlation_content, True),
    ("3_feature_vs_popularity.png", "FEATURE VS POPULARITY ANALİZİ",
     ["Energy ve danceability popularity ile pozitif trend", "Acousticness düşük popularity ile ilişkili", "Valence (mood) orta seviye popularity ile ilişkili"], False),
    ("4_genre_distribution.png", "GENRE DAĞILIMI", genre_content, True),
    ("5_year_distribution.png", "YIL BAZLI DAĞILIM",
     ["Veri seti 1960-2024 arası şarkıları kapsıyor", "2010 sonrası veri yoğunluğu artıyor"], False),
    ("6_feature_evolution.png", "AUDIO FEATURES EVOLUTION", evolution_content, True),
    ("7_boxplots.png", "OUTLIER ANALİZİ - BOX PLOTS",
     ["Loudness ve tempo'da önemli outlier'lar var", "Speechiness ve acousticness geniş dağılım gösteriyor", "Outlier'lar model eğitimi öncesi temizlendi"], False)
]:
    img_path = eda_dir / img_file
    if img_path.exists():
        if is_two_col:
            add_two_column_image_slide(prs, title, str(img_path), desc)
        else:
            add_image_slide(prs, title, str(img_path), desc)

# 15. Loudness Filtering Analysis
add_content_slide(prs, "LOUDNESS FİLTRELEME: VERİ TEMİZLEME", [
    "Spotify loudness ölçümüne göre -40 dB'nin altındaki kayıtlar filtrelendi",
    "Bu kayıtlar pratikte 'sessiz' veya sinyal içeriği zayıf kabul edilir",
    "Popülerlik modellemesi için uygun değil",
    "141,190 şarkıdan 12 tanesi filtrelendi (%0.01)",
    "Kalan 141,178 şarkı ile model eğitimi yapıldı",
    "Filtrelenen şarkıların loudness aralığı: -46.77 dB ile -40.05 dB"
])

img_path = eda_dir / "10_loudness_filtering_analysis.png"
if img_path.exists():
    add_two_column_image_slide(prs, "LOUDNESS FİLTRELEME ANALİZİ", str(img_path), [
        "Loudness dağılımı: Before vs After filtering",
        "Filtrelenen şarkıların loudness aralığı analizi",
        "Popularity karşılaştırması: Filtrelenen şarkılar düşük popularity",
        "Filtreleme sonrası veri kalitesi arttı",
        "Model eğitimi için daha temiz veri seti",
        "Threshold: -40 dB (müzikal olarak sessiz kabul edilen limit)"
    ])

img_path = eda_dir / "11_loudness_impact_analysis.png"
if img_path.exists():
    add_image_slide(prs, "LOUDNESS FİLTRELEME ETKİ ANALİZİ", str(img_path), [
        "Loudness dağılımı ve threshold çizgisi",
        "Loudness vs Popularity scatter plot",
        "Şarkı sayısı karşılaştırması: Before/After/Removed"
    ])

# 16. EDA Insights
add_content_slide(prs, "EDA İÇGÖRÜLERİ VE YORUMLAR", [
    "Energy ve danceability hit potansiyeli için kritik feature'lar",
    "Loudness modern müzikte artış trendi gösteriyor",
    "Loudness < -40 dB şarkılar filtrelendi (veri kalitesi için)",
    "Genre bazlı farklılıklar market-specific tahminler için önemli",
    "Temporal trends (yıl bazlı) model feature'larına dahil edildi",
    "Correlation analizi multi-collinearity tespiti için kullanıldı",
    "Outlier temizleme model performansını iyileştirdi",
    "Feature engineering: rare label encoding, scaling, normalization"
])

# 16-17. SPEECHINESS ANALYSIS (NEW)
add_content_slide(prs, "SPEECHINESS ANALİZİ: OUTLIER DEĞİL, TÜR ÖZELLİĞİ", [
    "Yüksek speechiness değerleri (>= 0.90) outlier değil, tür özelliği",
    "Rap ve hip-hop türlerinde yüksek speechiness normal",
    "Metal türünde de konuşma ağırlıklı şarkılar mevcut",
    "Speechiness: Müzikal melodiden ziyade konuşma/söz ağırlığını ölçer",
    "141,190 şarkıdan 158 tanesi speechiness >= 0.90 (%0.11)",
    "Bu şarkıların %42'si rap/hip-hop türünde"
])

img_path = eda_dir / "8_speechiness_genre_analysis.png"
if img_path.exists():
    add_two_column_image_slide(prs, "SPEECHINESS VS GENRE ANALİZİ", str(img_path), [
        "Rap ve hip-hop en yüksek ortalama speechiness'e sahip",
        "Rap: 0.214, Hip-hop: 0.221 ortalama speechiness",
        "Pop ve alternative düşük speechiness (0.063-0.064)",
        "Yüksek speechiness şarkılar rap/hip-hop'ta yoğunlaşıyor",
        "Speechiness >= 0.70 olan şarkıların %60'ı rap/hip-hop",
        "Bu durum tür özelliği, veri hatası değil"
    ])

img_path = eda_dir / "9_very_high_speechiness_analysis.png"
if img_path.exists():
    add_image_slide(prs, "ÇOK YÜKSEK SPEECHINESS ANALİZİ (>= 0.90)", str(img_path), [
        "158 şarkı speechiness >= 0.90 (toplamın %0.11'i)",
        "Rap: 46 şarkı, Hip-hop: 21 şarkı, Pop: 10 şarkı",
        "Yüksek speechiness şarkılar normal popularity seviyelerinde",
        "Bu değerler outlier değil, türün doğal özelliği"
    ])

add_content_slide(prs, "SPEECHINESS ANALİZİ SONUÇLARI", [
    "Yüksek speechiness değerleri (0.90+) outlier DEĞİL",
    "Rap ve hip-hop türlerinin doğal özelliği",
    "Metal türünde de konuşma ağırlıklı şarkılar normal",
    "Model eğitimi sırasında bu değerler korundu (temizlenmedi)",
    "Tür bazlı feature engineering ile daha iyi tahminler",
    "Speechiness, tür tanıma için önemli bir feature"
])

# ============================================================================
# MODEL PERFORMANCE SECTION
# ============================================================================

# Model Development Overview
add_content_slide(prs, "MODEL GELİŞTİRME SÜRECİ", [
    "9 farklı base model ile başlangıç performansı değerlendirildi",
    "XGBoost ve LightGBM en iyi performans gösteren modeller",
    "Grid Search CV ile hiperparametre optimizasyonu yapıldı",
    "Optimize edilmiş modeller Voting Regressor ile birleştirildi",
    "Final model: Voting Regressor (XGBoost + LightGBM ensemble)",
    "Cross-validation (CV=3) ile model performansı doğrulandı"
])

# Base Models Comparison
img_path = eda_dir / "12_base_models_comparison.png"
if img_path.exists():
    add_two_column_image_slide(prs, "BASE MODELLER KARŞILAŞTIRMASI", str(img_path), [
        "9 farklı regression modeli test edildi",
        "LightGBM: 16.5 RMSE (en iyi base model)",
        "XGBoost: 16.8 RMSE",
        "Random Forest: 18.5 RMSE",
        "Gradient Boosting: 17.2 RMSE",
        "Linear modeller daha yüksek RMSE gösterdi"
    ])

# Hyperparameter Optimization Impact
img_path = eda_dir / "13_hyperparameter_optimization_impact.png"
if img_path.exists():
    add_image_slide(prs, "HİPERPARAMETRE OPTİMİZASYONU ETKİSİ", str(img_path), [
        "XGBoost: 16.8 → 14.2 RMSE (%15.5 iyileşme)",
        "LightGBM: 16.5 → 14.0 RMSE (%15.2 iyileşme)",
        "Grid Search CV ile optimal hiperparametreler bulundu",
        "Her iki modelde de önemli performans artışı"
    ])

# Model Selection Pipeline
img_path = eda_dir / "14_model_selection_pipeline.png"
if img_path.exists():
    add_two_column_image_slide(prs, "MODEL SEÇİM SÜRECİ", str(img_path), [
        "Base Models: 25.5 RMSE (ortalama)",
        "Top Performers: 16.5 RMSE (XGBoost, LightGBM)",
        "Hyperparameter Optimization: 14.0 RMSE",
        "Voting Regressor: 13.8 RMSE",
        "Final Model: 13.8 RMSE (Test)",
        "Toplam iyileşme: %45.9 RMSE azalması"
    ])

# Final Model Performance Metrics
img_path = eda_dir / "15_final_model_performance_metrics.png"
if img_path.exists():
    add_image_slide(prs, "FİNAL MODEL PERFORMANS METRİKLERİ", str(img_path), None)

# Hyperparameter Optimization Details
img_path = eda_dir / "16_hyperparameter_optimization_details.png"
if img_path.exists():
    add_image_slide(prs, "HİPERPARAMETRE OPTİMİZASYONU DETAYLARI", str(img_path), [
        "XGBoost best params: n_estimators=900, learning_rate=0.03, max_depth=8",
        "LightGBM best params: n_estimators=1000, learning_rate=0.04, num_leaves=52",
        "Her iki model için regularization parametreleri optimize edildi",
        "Optimal hiperparametreler cross-validation ile doğrulandı"
    ])

# Model Performance Summary
add_content_slide(prs, "MODEL PERFORMANS ÖZETİ", [
    "Final Model: Voting Regressor (XGBoost + LightGBM)",
    "Test RMSE: 13.8",
    "Train RMSE: 12.5",
    "Overfitting Gap: 1.3 (düşük, iyi genelleme)",
    "R² Score: 0.78 (%78 varyans açıklanıyor)",
    "MAE: 9.2",
    "Model robust ve genelleme yapabiliyor"
])

# ============================================================================
# POPULARITY vs HIT SCORE ANALYSIS
# ============================================================================

# How Popularity and Hit Score are Calculated (Single Page)
img_path = eda_dir / "24_how_popularity_and_hitscore_calculated.png"
if img_path.exists():
    add_image_slide(prs, "POPULARITY vs HIT SCORE: NASIL HESAPLANIR?", str(img_path), [
        "POPULARITY: Kullanıcı davranışları → Spotify algoritması → Popularity skoru (0-100)",
        "HIT SCORE: Audio features → ML model → Percentile mapping → Hit Score (0-100)",
        "Popularity gerçek metrik (descriptive), Hit Score tahminsel metrik (predictive)",
        "İkisi farklı kaynaklardan ve yöntemlerle hesaplanır!"
    ])

# Hit Score Definition
add_two_column_slide(prs, "HIT SCORE: GİRDİ VE ÇIKTI",
    ["GİRDİLER (Inputs)", "• Model Prediction", "  - Voting Regressor output", "  - XGBoost + LightGBM", "  - Raw prediction value", "", "• Dataset Percentiles", "  - P10, P25, P50, P75", "  - P90, P95, P99", "  - Training data'dan hesaplanır"],
    ["HESAPLAMA (Calculation)", "1. Percentile pozisyonu belirle", "2. Percentile aralığını bul", "3. Linear interpolation uygula", "4. 0-100 skala'ya normalize et", "", "ÇIKTI (Output)", "• HIT SCORE", "  - Range: 0-100", "  - Normalize edilmiş skor", "", "ÖRNEK:", "Hit Score = 82.4", "(Yüksek hit potansiyeli)"]
)

# Transformation Flow
img_path = eda_dir / "22_popularity_to_hitscore_flow.png"
if img_path.exists():
    add_image_slide(prs, "POPULARITY → HIT SCORE DÖNÜŞÜM SÜRECİ", str(img_path), [
        "Popularity audio features ve user behavior'dan hesaplanır",
        "Model prediction, percentile-based mapping ile normalize edilir",
        "Hit Score, popularity'nin 0-100 arası standart skala versiyonu",
        "Dönüşüm süreci: Percentile analizi → Mapping → Normalizasyon"
    ])

# Mapping Table
img_path = eda_dir / "23_hitscore_mapping_table.png"
if img_path.exists():
    add_two_column_image_slide(prs, "PERCENTILE-BASED MAPPING TABLOSU", str(img_path), [
        "P99+ → 95-100: Very High Hit Potential",
        "P95-P99 → 90-95: High Hit Potential",
        "P75-P90 → 65-80: Good Hit Potential",
        "P50-P75 → 50-65: Moderate Hit Potential",
        "Linear interpolation ile kesin değer hesaplanır",
        "Bu mapping, farklı dataset'lerde karşılaştırılabilir skorlar üretir"
    ])

# Popularity vs Hit Score Overview
add_content_slide(prs, "POPULARITY vs HIT SCORE İLİŞKİSİ", [
    "Popularity: Spotify'ın gerçek popülerlik metriği (0-100)",
    "Hit Score: Model tahmininden türetilen normalize edilmiş skor (0-100)",
    "Percentile-based mapping ile popularity → hit score dönüşümü",
    "Yüksek correlation (r=0.997): Strong positive relationship",
    "Hit score, popularity'yi 0-100 arası standart skala'ya normalize eder",
    "Model predictions, percentile pozisyonuna göre hit score'a dönüştürülür"
])

# Popularity vs Hit Score Analysis
img_path = eda_dir / "20_popularity_vs_hit_score_analysis.png"
if img_path.exists():
    add_image_slide(prs, "POPULARITY vs HIT SCORE ANALİZİ", str(img_path), [
        "Scatter plot: Strong positive correlation",
        "Distribution comparison: Similar patterns",
        "Box plots: Hit score by popularity categories",
        "Correlation matrix: Audio features ile ilişkiler"
    ])

# Popularity vs Hit Score Relationship Details
img_path = eda_dir / "21_popularity_hit_relationship_details.png"
if img_path.exists():
    add_two_column_image_slide(prs, "POPULARITY vs HIT SCORE: DETAYLI İLİŞKİ ANALİZİ", str(img_path), [
        "Hexbin density plot: Popularity-Hit Score yoğunluk haritası",
        "Percentile analysis: Hit score percentiles by popularity",
        "Pearson correlation: 0.997 (very strong)",
        "Spearman correlation: 0.997 (monotonic relationship)",
        "Hit score, popularity'nin normalize edilmiş versiyonu",
        "Model predictions percentile-based mapping ile hit score'a dönüştürülüyor"
    ])

# Continue with remaining slides (18-28)
add_content_slide(prs, "SPOTIFY HIT PREDICTION", [
    "Real-time audio feature analizi (danceability, energy, valence, loudness, tempo)",
    "Interactive playlist ile şarkı seçimi",
    "Audio feature'lara göre benzer şarkı önerileri",
    "Market-specific hit tahminleri (6 farklı market)",
    "Adjustable feature sliders ile custom analiz",
    "Chart.js ile live görselleştirme",
    "Voting Ensemble Model (XGBoost + LightGBM + Random Forest + GBM)"
])

add_two_column_slide(prs, "INSTAGRAM PREDICTION ÖZELLİKLERİ",
    ["INSTAGRAM PICTURE", "• Image upload ve analiz", "• Visual feature extraction", "• Hit probability tahmini", "• Performance önerileri", "", "INSTAGRAM CAPTION", "• Text analiz ve sentiment", "• Hashtag optimizasyonu", "• Engagement önerileri"],
    ["FINAL PREDICTION", "• Combined score hesaplama", "• Weighted breakdown", "  - Spotify: 40%", "  - Picture: 30%", "  - Caption: 30%", "• Actionable öneriler", "", "REELS CONTENT", "• Multi-modal analiz", "• Entegre sistem"]
)

add_content_slide(prs, "VERİ VE METODOLOJİ", [
    "Spotify Dataset: 141K+ şarkı, audio features, popularity metrics",
    "Feature Engineering: Categorical encoding, rare label handling, scaling",
    "Preprocessing Pipeline: Automated data cleaning ve transformation",
    "Model Selection: Voting Regressor (ensemble of 4 models)",
    "Evaluation: RMSE, R² Score, cross-validation",
    "Market Adjustment: Base factors for 6 different markets",
    "Similarity Matching: Normalized Euclidean distance"
])

add_content_slide(prs, "MODEL PERFORMANSI", [
    "Ensemble Voting Model ile robust tahminler",
    "4 farklı algoritmanın kombinasyonu",
    "Feature importance analizi ile önemli özellikler belirlendi",
    "Market-specific tuning ile lokalize tahminler",
    "Real-time inference için optimize edilmiş pipeline",
    "Production-ready model deployment"
])

add_content_slide(prs, "ÖZELLİKLER VE KULLANICI DENEYİMİ", [
    "Modern, responsive web arayüzü (EA Games/Crytek tarzı)",
    "Interactive playlist ile kolay şarkı seçimi",
    "Real-time prediction updates",
    "Live görselleştirmeler (Chart.js)",
    "Market-specific tahmin seçenekleri",
    "Similar songs önerisi ile keşif",
    "Glassmorphism ve gradient efektler"
])

# ============================================================================
# UI/UX SECTION
# ============================================================================

# Homepage UI
img_path = eda_dir / "17_ui_homepage_mockup.png"
if img_path.exists():
    add_two_column_image_slide(prs, "WEB ARAYÜZÜ - ANA SAYFA", str(img_path), [
        "Modern dark theme tasarım (Spotify green accent)",
        "Hero section: Instagram Reels Hit Prediction",
        "Dynamic statistics: 141K+ tracks, 6 markets, 4 ML models",
        "Core systems showcase: Audio, Visual, Text analysis",
        "Glassmorphism navigation bar",
        "Responsive design for all devices"
    ])

# Spotify Prediction Page UI
img_path = eda_dir / "18_ui_spotify_prediction_mockup.png"
if img_path.exists():
    add_two_column_image_slide(prs, "WEB ARAYÜZÜ - SPOTIFY PREDICTION", str(img_path), [
        "Interactive playlist search and selection",
        "Real-time hit score prediction (0-100)",
        "Audio features visualization (Chart.js)",
        "Adjustable feature sliders for custom analysis",
        "Market-specific predictions (US, GB, DE, FR, JP, BR)",
        "Selected track information display",
        "Similar songs recommendation"
    ])

# Final Prediction Page UI
img_path = eda_dir / "19_ui_final_prediction_mockup.png"
if img_path.exists():
    add_two_column_image_slide(prs, "WEB ARAYÜZÜ - FINAL PREDICTION", str(img_path), [
        "Final combined hit score (weighted ensemble)",
        "Component breakdown: Spotify (40%), Picture (30%), Caption (30%)",
        "Visual score breakdown chart",
        "Actionable recommendations and insights",
        "Multi-modal integration display",
        "Score interpretation and suggestions"
    ])

add_content_slide(prs, "SONUÇLAR VE BAŞARILAR", [
    "Tam fonksiyonel Spotify hit prediction sistemi",
    "Production-ready Django web uygulaması",
    "Modüler mimari ile genişletilebilir yapı",
    "Docker ile kolay deployment",
    "Comprehensive API endpoints",
    "User-friendly interface ile pratik kullanım",
    "AWS Free Tier uyumlu altyapı"
])

add_content_slide(prs, "ÖĞRENİLENLER", [
    "End-to-end ML pipeline geliştirme (EDA → Model → Deployment)",
    "Feature engineering ve preprocessing teknikleri",
    "Ensemble learning modellerinin uygulanması",
    "Django ile production-ready web uygulaması geliştirme",
    "API design ve RESTful service mimarisi",
    "Docker containerization ve deployment",
    "Real-time ML inference optimizasyonu"
])

add_content_slide(prs, "GELECEK PLANLARI", [
    "Instagram dataset entegrasyonu ile gerçek veri kullanımı",
    "CNN-based image feature extraction",
    "NLP modelleri ile caption optimization (BERT, transformer models)",
    "User accounts ve prediction history tracking",
    "A/B testing önerileri",
    "Export predictions (CSV/PDF)",
    "Mobile app geliştirme",
    "Advanced analytics dashboard"
])

# Teşekkürler
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = DARK_GRAY

left = Inches(0.5)
top = Inches(2)
width = Inches(9)
height = Inches(1.5)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.text = "TEŞEKKÜRLER"
tf.paragraphs[0].font.size = Pt(54)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = SPOTIFY_GREEN
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

top = Inches(4)
height = Inches(2)
txBox2 = slide.shapes.add_textbox(left, top, width, height)
tf2 = txBox2.text_frame
tf2.text = "Sorularınız için hazırım!\n\nEnes TOP\nMIUUL 19. Dönem Data Scientist Bootcamp\nAralık 2025"
tf2.paragraphs[0].font.size = Pt(24)
tf2.paragraphs[0].font.color.rgb = WHITE
tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
for para in tf2.paragraphs[1:]:
    para.font.size = Pt(20)
    para.font.color.rgb = WHITE
    para.alignment = PP_ALIGN.CENTER

# Save
import datetime
timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
output_file = f"MIUUL_Final_Sunum_Complete_Final_{timestamp}.pptx"
prs.save(output_file)
print(f"Sunum başarıyla oluşturuldu: {output_file}")
print(f"Toplam slayt sayısı: {len(prs.slides)}")

